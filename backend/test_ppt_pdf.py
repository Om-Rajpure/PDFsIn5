import io
from fastapi.testclient import TestClient
from pptx import Presentation
from app.main import app

def test_ppt_to_pdf():
    client = TestClient(app)
    
    # Create a dummy PowerPoint file in memory
    ppt_io = io.BytesIO()
    
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PDFsIn5 PowerPoint to PDF Test"
        subtitle.text = "This is a dynamically generated PPT file for backend testing."
        
        prs.save(ppt_io)
    except ImportError:
        print("python-pptx is not installed, cannot generate dynamic PPT file. Skipping test.")
        return
        
    ppt_io.seek(0)
    
    print("Sending request to /api/tools/ppt-to-pdf...")
    response = client.post(
        "/api/tools/ppt-to-pdf",
        files={"files": ("test.pptx", ppt_io, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
    )
    
    print(f"Status Code: {response.status_code}")
    if response.status_code == 200:
        print("Success! PDF received.")
        print(f"Headers: {response.headers}")
        print(f"Content Length: {len(response.content)} bytes")
        with open("ppt_test_output.pdf", "wb") as f:
            f.write(response.content)
        print("Saved ppt_test_output.pdf for manual inspection.")
    else:
        print(f"Error: {response.json()}")

if __name__ == "__main__":
    test_ppt_to_pdf()
